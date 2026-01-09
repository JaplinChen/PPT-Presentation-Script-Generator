"""
PPT Audio/Video Embedding Service.

This module handles the injection of generated audio and video assets into PowerPoint slides.
Refactored for enterprise standards:
- Strong typing
- Clear separation of concerns
- Robust error handling
- Constants for configuration
"""

import logging
import os
import re
import shutil
import uuid
import collections
import collections.abc
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any, Union

import cv2
from mutagen.mp3 import MP3
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.oxml import parse_xml # type: ignore
from pptx.oxml.ns import qn     # type: ignore

# --- Monkey Patch for python-pptx ---
# Fixes "AttributeError: 'Part' object has no attribute 'sha1'"
import pptx.opc.package # type: ignore

# Fixes for python 3.10+ container imports
if not hasattr(collections, 'Container'):
    collections.Container = collections.abc.Container # type: ignore
if not hasattr(collections, 'Mapping'):
    collections.Mapping = collections.abc.Mapping # type: ignore
if not hasattr(collections, 'MutableMapping'):
    collections.MutableMapping = collections.abc.MutableMapping # type: ignore

if not hasattr(pptx.opc.package.Part, 'sha1'):
    @property
    def sha1_patch(self):
        return f"patched_sha1_{id(self)}"
    pptx.opc.package.Part.sha1 = sha1_patch

logger = logging.getLogger(__name__)

# --- Constants ---
SLIDE_1_VIDEO_SCALE_FACTOR = 0.15          # 15% of ORIGINAL video size
SLIDE_1_PHOTO_SCALE_FACTOR = 0.15          # 15% of SLIDE width (Legacy logic preserved)
SLIDE_OTHER_MEDIA_SIZE_CM = 2.0            # 2.0 cm compacted size
MARGIN_DEFAULT_CM = 0.5
MARGIN_COVER_CM = 1.0
AUDIO_ICON_SIZE_CM = 0.85
EMU_PER_PIXEL_96DPI = 9525                 # 1 pixel = 9525 EMU (approx)

class PPTEmbedder:
    """
    Handles embedding of audio and video into PowerPoint presentations.
    """

    def __init__(self, output_dir: Path):
        self.output_dir = output_dir

    async def embed_both(
        self,
        original_pptx_path: str,
        slide_scripts: List[Dict[str, Any]],
        audio_paths: List[Optional[str]],
        video_paths: List[Optional[str]],
        photo_path: Optional[str] = None,
        progress_callback: Optional[callable] = None,
    ) -> str:
        """
        Embeds both audio and video into the PPT.
        Priority: Video > Static Photo + Audio > Only Audio
        """
        original_path = Path(original_pptx_path)
        output_filename = self._get_sequential_filename(original_path, "final")
        output_path = self.output_dir / output_filename
        
        # Copy original first
        shutil.copy(original_path, output_path)

        prs = Presentation(output_path)
        total_slides = len(prs.slides)
        
        # Filter visible slides
        visible_slides = [s for s in prs.slides if str(s.element.get('show')) not in ('0', 'false')]
        
        for idx, slide in enumerate(visible_slides):
            slide_no = idx + 1
            if progress_callback:
                progress_callback(
                    int((slide_no / total_slides) * 100), 
                    f"Assembling slide {slide_no}/{total_slides}..."
                )

            # Get Assets
            video_path = self._resolve_path(video_paths[idx]) if idx < len(video_paths) else None
            audio_path = self._resolve_path(audio_paths[idx]) if idx < len(audio_paths) else None
            
            try:
                # Strategy 1: Video
                if video_path and os.path.exists(video_path):
                    self._embed_video_strategy(slide, video_path, prs, slide_no, photo_path)
                    
                # Strategy 2: Static Photo + Audio
                elif photo_path and os.path.exists(photo_path) and audio_path and os.path.exists(audio_path):
                     # Validate audio
                    duration = self._get_audio_duration(audio_path)
                    self._embed_static_avatar_strategy(slide, photo_path, audio_path, prs, duration, is_title=(idx==0))
                    
                # Strategy 3: Audio Only
                elif audio_path and os.path.exists(audio_path):
                    duration = self._get_audio_duration(audio_path)
                    self._embed_audio_only_strategy(slide, audio_path, prs, duration)
                    
            except Exception as e:
                logger.error(f"Failed to process media for slide {slide_no}: {e}", exc_info=True)

        prs.save(output_path)
        return str(output_path.resolve())

    async def embed_audio(
        self,
        original_pptx_path: str,
        slide_scripts: List[Dict[str, Any]],
        audio_generator: Any,
        voice: str,
        rate: str = "+0%",
        pitch: str = "+0Hz",
        progress_callback: Optional[callable] = None,
    ) -> Tuple[str, Dict[int, str], List[str]]:
        """
        Generates and embeds audio based on scripts.
        """
        original_path = Path(original_pptx_path)
        output_filename = self._get_sequential_filename(original_path)
        output_path = self.output_dir / output_filename
        shutil.copy(original_path, output_path)
        
        prs = Presentation(output_path)
        all_slide_scripts: Dict[int, str] = {}
        audio_files: List[str] = []
        
        script_data_map = {int(item['slide_no']): item for item in slide_scripts}
        visible_slides = [s for s in prs.slides if str(s.element.get('show')) not in ('0', 'false')]
        total = len(visible_slides)

        for idx, slide in enumerate(visible_slides):
            slide_no = idx + 1
            if progress_callback:
                progress_callback(int((slide_no/total)*80), f"Processing slide {slide_no}/{total}...")

            data = script_data_map.get(slide_no)
            if not data:
                continue

            raw_script = data.get('script', '')
            all_slide_scripts[slide_no] = raw_script
            
            clean_script = self._clean_script_text(raw_script)
            if not clean_script:
                continue

            try:
                audio_info = await audio_generator.generate_audio(
                    clean_script, voice, rate, pitch, filename=f"slide_{slide_no}.mp3"
                )
                audio_path = audio_info['path']
                audio_files.append(audio_path)
                
                duration = self._get_audio_duration(audio_path)
                self._embed_audio_only_strategy(slide, audio_path, prs, duration)
                
            except Exception as e:
                logger.error(f"Generate/Embed audio failed for slide {slide_no}: {e}")

        prs.save(output_path)
        return str(output_path.resolve()), all_slide_scripts, audio_files

    async def embed_videos(
        self,
        pptx_path: str,
        video_paths: List[Optional[str]],
        progress_callback: Optional[callable] = None,
    ) -> str:
        """
        Embeds pre-generated videos into an existing PPTX.
        """
        prs = Presentation(pptx_path)
        total = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            if i >= len(video_paths) or not video_paths[i]:
                continue
            
            v_path = video_paths[i] # type: ignore
            if not v_path or not os.path.exists(v_path):
                continue

            if progress_callback:
                progress_callback(int((i/total)*100), f"Embedding video slide {i+1}...")

            try:
                # Just using Strategy 1
                self._embed_video_strategy(slide, v_path, prs, slide_no=i+1)
            except Exception as e:
                logger.error(f"Embed video failed slide {i+1}: {e}")

        prs.save(pptx_path)
        return pptx_path

    # --- Strategy Implementations ---

    def _embed_video_strategy(self, slide, video_path: str, prs, slide_no: int, poster_path: Optional[str] = None):
        """
        Strategy 1: Embed Video
        Slide 1: 15% of ORIGINAL VIDEO SIZE (Pixel-based). Bottom Right.
        Slide 2+: 2.0cm width. Bottom Left.
        """
        slide_w, slide_h = prs.slide_width, prs.slide_height
        
        # 1. Measurement
        orig_w, orig_h = self._get_video_dimensions_robust(video_path)
        aspect_ratio = orig_w / orig_h if orig_h > 0 else 1.0

        # 2. Sizing Logic
        if slide_no == 1:
            # FIX: Scale Width to 15% of SLIDE WIDTH
            # Calculate height based on Aspect Ratio to prevent distortion
            v_width = int(slide_w * 0.15)
            v_height = int(v_width / aspect_ratio)
            logger.info(f"Slide 1 Video Scaling: 15% Width = {v_width} EMU, Ratio={aspect_ratio:.2f}")
        else:
            # Slide 2+: 5% of SLIDE WIDTH
            v_width = int(slide_w * 0.05)
            v_height = int(v_width / aspect_ratio)

        # Ensure aspect ratio is preserved
        # v_height = int(v_width / aspect_ratio) # (This logic is now integrated above)

        # 3. Positioning Logic
        if slide_no == 1:
            margin = Cm(MARGIN_COVER_CM)
            left = slide_w - v_width - margin
            top = slide_h - v_height - margin
        else:
            margin = Cm(MARGIN_DEFAULT_CM)
            left = margin
            top = slide_h - v_height - margin

        # 4. Insertion
        mime = self._detect_mime(video_path)
        movie = slide.shapes.add_movie(
            video_path, left=left, top=top, width=v_width, height=v_height,
            poster_frame_image=poster_path, mime_type=mime
        )

        # 5. Styling (Oval Mask)
        # Always apply oval mask as requested
        self._set_shape_to_oval(movie)

        # 6. Autoplay
        self._add_autoplay_timing(slide, movie.shape_id)

    def _embed_static_avatar_strategy(self, slide, photo_path: str, audio_path: str, prs, duration: float, is_title: bool):
        """
        Strategy 2: Static Photo + Audio Icon
        Slide 1: Photo Bottom Right (15% Slide Width). Audio hidden/small.
        Slide 2+: Photo Bottom Left (0.85cm). Audio Icon Next to it.
        """
        slide_w, slide_h = prs.slide_width, prs.slide_height
        
        if is_title:
            # === Title Slide ===
            # Sizing: 15% of SLIDE WIDTH (Legacy requirement for photos)
            target_w = slide_w * SLIDE_1_PHOTO_SCALE_FACTOR
            
            # Add at 0,0 to get height (maintain aspect ratio)
            pic = slide.shapes.add_picture(photo_path, 0, 0, width=target_w)
            
            # Position: Bottom Right
            margin = Cm(MARGIN_COVER_CM)
            pic.left = int(slide_w - target_w - margin)
            pic.top = int(slide_h - pic.height - margin)
            
            # Audio: Discreet/Invisible
            icon_sz = Cm(AUDIO_ICON_SIZE_CM)
            self._insert_audio_shape(
                slide, audio_path, 
                left=slide_w - icon_sz - Cm(0.5),
                top=slide_h - icon_sz - Cm(0.5), # Behind/near photo
                width=icon_sz, height=icon_sz
            )
            
        else:
            # === Content Slide ===
            # Sizing: Compact (0.85cm)
            size = Cm(AUDIO_ICON_SIZE_CM)
            margin = Cm(MARGIN_DEFAULT_CM)
            
            # Photo
            pic = slide.shapes.add_picture(photo_path, 0, 0, width=size)
            pic.left = int(margin)
            pic.top = int(slide_h - margin - pic.height)
            
            img_right = pic.left + pic.width
            
            # Audio Icon
            self._insert_audio_shape(
                slide, audio_path,
                left=img_right + Cm(0.1),
                top=slide_h - margin - size,
                width=size, height=size,
                show_poster=True
            )

        # Transition Logic
        self._set_slide_transition(slide, duration)

    def _embed_audio_only_strategy(self, slide, audio_path: str, prs, duration: float):
        """
        Strategy 3: Audio Only
        Icon Bottom Left.
        """
        icon_sz = Cm(0.5)
        margin = Cm(0.5)
        
        self._insert_audio_shape(
            slide, audio_path,
            left=margin,
            top=prs.slide_height - icon_sz - margin,
            width=icon_sz, height=icon_sz,
            show_poster=True
        )
        self._set_slide_transition(slide, duration)

    # --- Helpers ---

    def _resolve_path(self, path: Optional[str]) -> Optional[str]:
        """Resolves relative /outputs/ paths to absolute."""
        if not path:
            return None
        if path.startswith('/outputs/'):
            return str(self.output_dir / path[len('/outputs/'):])
        return path

    def _get_audio_duration(self, path: str) -> float:
        try:
            return MP3(path).info.length
        except:
            return 0.0

    def _clean_script_text(self, text: str) -> str:
        """Cleans script text for TTS."""
        if not text: return ""
        text = re.sub(r'===.*?===', '', text)
        text = re.sub(r'---.*?---', '', text)
        text = re.sub(r'\([約大概]*\s*\d+\s*[秒分鐘seconds]+\)', '', text)
        text = re.sub(r'【.*?預計時間.*?】', '', text)
        # Custom
        text = text.replace("VPIC1", "VPIC one")
        text = re.sub(r'[*()[]/]', ' ', text)
        return re.sub(r'\s+', ' ', text).strip()

    def _get_sequential_filename(self, original_path: Path, prefix: str = None) -> str:
        """Generates OriginalName_YYYYMMDD_001.pptx"""
        stem = original_path.stem
        # Remove existing suffix
        stem = re.sub(r'_\d{8}_\d{3}$', '', stem)
        today = datetime.now().strftime("%Y%m%d")
        base = f"{stem}_{today}"
        
        max_seq = 0
        for f in self.output_dir.glob(f"{base}_*.pptx"):
            try:
                # Expected: ..._001.pptx
                parts = f.stem.split('_')
                if parts and parts[-1].isdigit() and len(parts[-1]) == 3:
                     seq = int(parts[-1])
                     if seq > max_seq: max_seq = seq
            except: pass
            
        return f"{base}_{max_seq+1:03d}.pptx"

    def _get_video_dimensions_robust(self, video_path: str) -> Tuple[int, int]:
        """
        Robustly gets video dimensions using OpenCV.
        Handles Unicode paths by copying to a temp file if needed.
        """
        def _read(p):
            cap = cv2.VideoCapture(p)
            if cap.isOpened():
                w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
                h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                cap.release()
                return w, h
            return 0, 0

        # Attempt 1: Direct
        try:
            w, h = _read(video_path)
            if w > 0: return w, h
        except Exception: 
            pass

        # Attempt 2: Temp copy (Unicode/Path fix)
        try:
            logger.info(f"Direct read failed for {video_path}, using temp workaround.")
            folder = os.path.dirname(video_path)
            ext = os.path.splitext(video_path)[1]
            temp_name = f"probe_{uuid.uuid4().hex[:8]}{ext}"
            temp_path = os.path.join(folder, temp_name)
            
            shutil.copy2(video_path, temp_path)
            w, h = _read(temp_path)
            
            try: os.remove(temp_path)
            except: pass
            
            if w > 0: return w, h
        except Exception as e:
            logger.warning(f"Video dim detect failed: {e}")
            
        return 0, 0

    def _detect_mime(self, path: str) -> str:
        lower = path.lower()
        if lower.endswith('.webm'): return 'video/webm'
        if lower.endswith('.mov'): return 'video/quicktime'
        return 'video/mp4'

    def _insert_audio_shape(self, slide, audio_path, left, top, width, height, show_poster=False):
        """Inserts audio movie shape."""
        poster = None
        if show_poster:
            asset = Path(__file__).parent.parent.parent / "assets" / "audio_icon.png"
            if asset.exists(): poster = str(asset)
        
        try:
            m = slide.shapes.add_movie(
                audio_path, left=left, top=top, width=width, height=height,
                poster_frame_image=poster, mime_type='audio/mp3'
            )
            self._add_autoplay_timing(slide, m.shape_id)
        except Exception:
            # Retry without poster
            m = slide.shapes.add_movie(
                audio_path, left=left, top=top, width=width, height=height,
                poster_frame_image=None, mime_type='audio/mp3'
            )
            self._add_autoplay_timing(slide, m.shape_id)

    def _set_slide_transition(self, slide, duration_sec: float):
        if hasattr(slide, 'slide_show_transition'):
            t = slide.slide_show_transition
            t.advance_on_time = True
            t.advance_after_time = int((duration_sec + 2.0) * 1000)

    # --- XML/OpenXML Manipulations ---

    def _add_autoplay_timing(self, slide, shape_id):
        """Injects complex animation XML to auto-play media."""
        timing = slide.element.find(qn('p:timing'))
        if timing is None:
            timing = parse_xml('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
            extLst = slide.element.find(qn('p:extLst'))
            if extLst is not None: slide.element.insert(slide.element.index(extLst), timing)
            else: slide.element.append(timing)
        
        tnLst = timing.find(qn('p:tnLst'))
        if tnLst is None:
            tnLst = parse_xml('<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>')
            timing.append(tnLst)
        
        # Clear existing to prevent dupes/conflicts
        for child in list(tnLst): tnLst.remove(child)
        
        # Standard Autosynced Animation XML
        xml = f"""
        <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                        <p:childTnLst>
                          <p:cmd type="call" cmd="playFrom(0.0)">
                            <p:cBhvr>
                              <p:cTn id="4" dur="1" fill="hold"/>
                              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                            </p:cBhvr>
                          </p:cmd>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
        """
        tnLst.append(parse_xml(xml))

    def _set_shape_to_oval(self, shape):
        """Applies Oval geometry to shape via XML."""
        try:
            from lxml import etree # type: ignore
            spPr = shape.element.find(qn('p:spPr'))
            if spPr is None:
                 # Video shapes (p:pic) usually have spPr as child
                 for child in shape.element:
                     if 'spPr' in child.tag:
                         spPr = child
                         break
            if spPr is None: return

            # Clear geometry/line
            for child in list(spPr):
                tag = etree.QName(child.tag).localname
                if tag in ('prstGeom', 'custGeom', 'ln'):
                    spPr.remove(child)
            
            # Add Ellipse
            a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            prstGeom = etree.SubElement(spPr, f'{{{a_ns}}}prstGeom')
            prstGeom.set('prst', 'ellipse')
            etree.SubElement(prstGeom, f'{{{a_ns}}}avLst')
            
        except ImportError:
            pass # lxml not available?
        except Exception as e:
            logger.debug(f"Oval mask failed: {e}")
