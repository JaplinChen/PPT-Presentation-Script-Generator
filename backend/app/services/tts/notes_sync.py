"""
PowerPoint Notes Synchronization Service.

Handles writing script text to slide notes using multiple fallback methods.
Cross-platform compatible (Windows/macOS).
"""
from pathlib import Path
from typing import Dict, Optional
import logging

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)


class NotesSync:
    """
    Synchronizes slide notes with script text.
    
    Uses multiple fallback methods to ensure notes are written:
    1. Standard notes_text_frame
    2. BODY placeholder
    3. Any writable text frame
    4. XML-level shape creation
    """
    
    def sync_notes(self, pptx_path: str, slide_scripts: Dict[int, str]) -> int:
        """
        Synchronize slide notes using python-pptx.
        
        Args:
            pptx_path: Absolute path to PPTX file
            slide_scripts: Dict mapping slide number (1-indexed) to script text
            
        Returns:
            Number of slides successfully updated
        """
        path = Path(pptx_path)
        if not path.exists():
            logger.error(f"[NotesSync] File not found: {pptx_path}")
            return 0

        try:
            logger.info(f"[NotesSync] Starting sync for {path.name}")
            prs = Presentation(str(path))
            
            # Get visible slides only
            visible_slides = [s for s in prs.slides if s.element.get('show') != '0']
            logger.debug(f"[NotesSync] {len(visible_slides)} visible slides, {len(slide_scripts)} scripts")
            
            synced_count = 0
            for slide_no, script_text in slide_scripts.items():
                if not script_text:
                    continue
                
                if not (1 <= slide_no <= len(visible_slides)):
                    logger.warning(f"[NotesSync] Slide {slide_no} out of range")
                    continue
                
                try:
                    slide = visible_slides[slide_no - 1]
                    notes_slide = slide.notes_slide
                    
                    if self._write_notes(notes_slide, script_text):
                        synced_count += 1
                        logger.debug(f"[NotesSync] Slide {slide_no}: synced")
                    else:
                        logger.debug(f"[NotesSync] Slide {slide_no}: Failed")
                        
                except Exception as e:
                    logger.error(f"[NotesSync] Slide {slide_no}: {e}")
            
            prs.save(str(path))
            logger.info(f"[NotesSync] ✅ Done - {synced_count}/{len(slide_scripts)} updated")
            return synced_count
            
        except Exception as e:
            logger.error(f"[NotesSync] Error: {e}", exc_info=True)
            return 0

    def _write_notes(self, notes_slide, text: str) -> bool:
        """
        Try multiple methods to write notes text.
        
        Args:
            notes_slide: The notes slide object
            text: The script text to write
            
        Returns:
            True if successful, False otherwise
        """
        # Method 1: Standard notes_text_frame - clear and rebuild paragraphs
        try:
            tf = notes_slide.notes_text_frame
            if tf:
                # Clear all existing paragraphs except the first one
                p = tf.paragraphs[0]
                p.clear()
                
                # Split text by newlines and add as runs/paragraphs
                lines = text.split('\n')
                for i, line in enumerate(lines):
                    if i == 0:
                        p.text = line
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                return True
        except Exception as e:
            logger.debug(f"[NotesSync] Method 1 failed: {e}")
        
        # Method 2: Find BODY placeholder (idx=2) specifically
        for shape in notes_slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                try:
                    phf = shape.placeholder_format
                    if phf.type == PP_PLACEHOLDER.BODY or phf.idx == 2:
                        tf = shape.text_frame
                        # Clear and write
                        for para in tf.paragraphs:
                            para.clear()
                        tf.paragraphs[0].text = text
                        return True
                except Exception as e:
                    logger.debug(f"[NotesSync] Method 2 failed: {e}")
        
        # Method 3: Find any writable text frame (skip slide image at idx=1)
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                try:
                    if hasattr(shape, 'placeholder_format'):
                        if shape.placeholder_format.idx == 1:
                            continue
                    shape.text_frame.paragraphs[0].text = text
                    return True
                except Exception:
                    pass
        
        # Method 4: Create text shape via XML
        return self._create_text_shape_xml(notes_slide, text)

    def _create_text_shape_xml(self, notes_slide, text: str) -> bool:
        """
        Create a text shape using direct XML manipulation.
        Dynamically positions the text box below the slide image if present.
        
        Args:
            notes_slide: The notes slide object
            text: The script text to write
            
        Returns:
            True if successful, False otherwise
        """
        try:
            sp_tree = notes_slide.shapes._spTree
            
            # Default Layout Metrics (EMUs)
            # 1 inch = 914400 EMUs
            PAGE_WIDTH = 6858000   # 7.5 inches
            PAGE_HEIGHT = 9144000  # 10 inches
            MARGIN_X = 685800      # 0.75 inch
            MARGIN_Y = 457200      # 0.5 inch
            
            # 1. Attempt to find the Slide Image (sys-idx 1 or known placeholder type) to span below it
            ref_y = 4572000 # Default to 5.0 inches (middle of page) if nothing found
            
            for shape in notes_slide.shapes:
                try:
                    # Check for Slide Image Placeholder
                    if shape.is_placeholder:
                        # PP_PLACEHOLDER.SLIDE_IMAGE is usually 1, or search by idx
                        if shape.placeholder_format.idx == 1:
                            # Found the slide image!
                            # Calculate simple bottom metric
                            # Note: python-pptx allows direct access to top/height in EMUs
                            bottom = shape.top + shape.height
                            ref_y = bottom + MARGIN_Y
                            logger.debug(f"[NotesSync] Found Slide Image: bottom={bottom}, set ref_y={ref_y}")
                            break
                except Exception:
                    pass
            
            # Ensure we don't go off-page
            if ref_y > PAGE_HEIGHT - 914400: # If remaining space < 1 inch
                ref_y = 5029200 # Fallback to ~5.5 inches
                logger.warning("[NotesSync] Dynamic ref_y too low, resetting to safe default")

            # Calculate remaining height with bottom margin
            avail_height = PAGE_HEIGHT - ref_y - MARGIN_Y
            if avail_height < 914400:
                 avail_height = 3657600 # Fallback 4 inches

            # Generate unique shape ID
            new_id = self._get_next_shape_id(sp_tree)
            
            # Create shape XML with text
            # CRITICAL: Added <p:ph type="body" idx="2"/> to make it a recognized placeholder
            sp_xml = f'''
            <p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                <p:nvSpPr>
                    <p:cNvPr id="{new_id}" name="Notes Placeholder {new_id}"/>
                    <p:cNvSpPr txBox="1"/>
                    <p:nvPr>
                        <p:ph type="body" idx="2"/>
                    </p:nvPr>
                </p:nvSpPr>
                <p:spPr>
                    <a:xfrm>
                        <a:off x="{MARGIN_X}" y="{int(ref_y)}"/>
                        <a:ext cx="{PAGE_WIDTH - (MARGIN_X * 2)}" cy="{int(avail_height)}"/>
                    </a:xfrm>
                    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
                </p:spPr>
                <p:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p><a:r><a:rPr lang="zh-TW" sz="1200"/><a:t></a:t></a:r></a:p>
                </p:txBody>
            </p:sp>
            '''
            
            sp_elem = etree.fromstring(sp_xml)
            
            # Set text content
            t_elem = sp_elem.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}t')
            if t_elem is not None:
                t_elem.text = text
            
            sp_tree.append(sp_elem)
            logger.info(f"[NotesSync] Created dynamic XML notes placeholder id={new_id} at y={ref_y}")
            return True
            
        except Exception as e:
            logger.error(f"[NotesSync] XML creation failed: {e}")
            return False

    def _get_next_shape_id(self, sp_tree) -> int:
        """Get the next available shape ID from the spTree."""
        max_id = 0
        for sp in sp_tree.iter(qn('p:sp')):
            nv_sp_pr = sp.find(qn('p:nvSpPr'))
            if nv_sp_pr is not None:
                c_nv_pr = nv_sp_pr.find(qn('p:cNvPr'))
                if c_nv_pr is not None:
                    try:
                        shape_id = int(c_nv_pr.get('id', 0))
                        max_id = max(max_id, shape_id)
                    except ValueError:
                        pass
        return max_id + 1
