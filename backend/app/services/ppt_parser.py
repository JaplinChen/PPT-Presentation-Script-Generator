from pptx import Presentation
from typing import List, Dict, Any
from pathlib import Path
import time
import logging

logger = logging.getLogger(__name__)

class PPTParser:
    """Parsing PowerPoint files and extracting structured content - Optimized Version"""
    
    def parse(self, ppt_path: str) -> List[Dict[str, Any]]:
        """
        Parse PPT and extract content from all slides.
        """
        start_time = time.time()
        
        if not Path(ppt_path).exists():
            raise FileNotFoundError(f"PPT file not found: {ppt_path}")
        
        try:
            # Presentation loading involves I/O and can be slow
            prs = Presentation(ppt_path)
        except Exception as e:
            logger.error(f"Failed to read PPT file: {e}")
            raise
            
        slides_data = []
        slide_height = prs.slide_height
        visible_slide_no = 0
        
        # Pre-calculated threshold
        threshold = slide_height * 0.25
        
        for i, slide in enumerate(prs.slides):
            # Check for hidden slides
            if slide.element.get('show') == '0' or slide.element.get('show') == 'false':
                continue
            
            visible_slide_no += 1
            
            title = ""
            bullets = []
            tables = []
            image_count = 0
            candidate_titles = []
            
            # Single pass through all shapes to reduce attribute access
            for shape in slide.shapes:
                try:
                    stype = shape.shape_type
                except Exception: continue
                
                # 1. Image count (Fast path)
                if stype == 13: # Picture
                    image_count += 1
                    continue
                
                # 2. Table extraction (19 = Table)
                if stype == 19 or shape.has_table:
                    try:
                        table = shape.table
                        table_data = {
                            'rows': len(table.rows),
                            'cols': len(table.columns),
                            'content': [[cell.text.strip() for cell in row.cells] for row in table.rows]
                        }
                        tables.append(table_data)
                    except Exception: pass
                    continue
                
                # 3. Text content extraction
                # Only for shapes with text frames (Placeholder=14, Textbox=17, AutoShape=1)
                if shape.has_text_frame:
                    try:
                        text = shape.text.strip()
                        if not text:
                            continue
                            
                        # Check title placeholder (Placeholder type 1 = Title)
                        is_title = False
                        if shape.is_placeholder:
                            ph_type = shape.placeholder_format.type
                            if ph_type == 1 or ph_type == 3: # 1=Title, 3=Centered Title
                                if not title: title = text
                                is_title = True
                        
                        if not is_title:
                            # Candidate titles (top of slide, larger font)
                            if hasattr(shape, "top") and shape.top < threshold and len(text) > 1:
                                fsize = 0
                                try:
                                    if shape.text_frame.paragraphs[0].runs:
                                        fsize = shape.text_frame.paragraphs[0].runs[0].font.size or 0
                                except Exception: pass
                                candidate_titles.append({"text": text, "top": shape.top, "fsize": fsize})
                            
                            # Content bullets (filter extremely short text)
                            if len(text) > 1:
                                lines = [l.strip() for l in text.split('\n') if l.strip()]
                                bullets.extend(lines)
                    except Exception: pass
            
            # Best effort title selection
            if not title and candidate_titles:
                # Sort by font size desc, then top position asc
                candidate_titles.sort(key=lambda x: (-x['fsize'], x['top']))
                title = candidate_titles[0]['text']
            
            # Notes extraction (lazy access)
            notes = ""
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception: pass
            
            slides_data.append({
                'slide_no': visible_slide_no,
                'title': title,
                'bullets': bullets,
                'tables': tables,
                'notes': notes,
                'image_count': image_count
            })
            
            # Log progress every 50 slides
            if visible_slide_no % 50 == 0:
                logger.info(f"Processed {visible_slide_no} slides...")
        
        elapsed = time.time() - start_time
        logger.info(f"Completed! Total {len(slides_data)} slides, Time: {elapsed:.2f}s")
        return slides_data
    
    def get_summary(self, slides_data: List[Dict[str, Any]]) -> Dict[str, int]:
        """Get PPT summary (O(N) traversal)"""
        return {
            'total_slides': len(slides_data),
            'titled_slides': sum(1 for s in slides_data if s['title']),
            'total_bullets': sum(len(s['bullets']) for s in slides_data),
            'slides_with_notes': sum(1 for s in slides_data if s['notes']),
            'slides_with_tables': sum(1 for s in slides_data if s['tables']),
            'total_images': sum(s['image_count'] for s in slides_data)
        }
